import os, io, re, traceback, tempfile, base64, shutil, copy
from itertools import cycle
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document as DocxDocument
from docx.shared import Pt as DocxPt, Cm, Inches, RGBColor as DocxRGB, Emu as DocxEmu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from lxml import etree

SLIDE_W, SLIDE_H = 12192000, 6858000
VINHO = RGBColor(0x70, 0x00, 0x1C)
PRETO = RGBColor(0x00, 0x00, 0x00)
LOGO_MASTER_X,LOGO_MASTER_Y,LOGO_MASTER_CX,LOGO_MASTER_CY = 9282544,385975,2507673,776504
LOGO_CAPA_X,LOGO_CAPA_Y,LOGO_CAPA_CX,LOGO_CAPA_CY = 4429838,558156,3251122,764208
LOGO_ENC_X,LOGO_ENC_Y,LOGO_ENC_CX,LOGO_ENC_CY = 2756357,2683952,6679286,1490095
TEXTO_X,TEXTO_Y,TEXTO_CX,TEXTO_H = 347272,1074509,11497455,5200000
RODAPE_X,RODAPE_Y,RODAPE_CX,RODAPE_CY = 3382781,6298579,8461946,400110
CAPA_X,CAPA_Y,CAPA_CX,CAPA_CY = 2548009,2013228,7465441,2985433
AREA_UTIL = (RODAPE_Y - TEXTO_Y) - int(40 * 12700)
FAIXA_H = 190000   # faixa vinho no topo (~15pt)
FAIXA_Y = 0        # colada no topo, sem espacamento
LARGURA = 11497455
FATOR = 0.42
CITACOES = [
    '"A forca esta em se erguer com mais poder, como a aguia."',
    '"A aguia foca no futuro, nao no que esta atras."',
    '"Para voar alto, deixe as correntes para tras, como a aguia."',
    '"Nas tempestades, como a aguia, encontre o voo mais alto."',
]
ASSETS = os.path.join(os.path.dirname(__file__), "assets")
MODELO_DOCX = os.path.join(os.path.dirname(__file__), "assets", "MODELO_DOCUMENTO.docx")
WS_NS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
NS_R  = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
NS_A  = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
VML_NS = '{urn:schemas-microsoft-com:vml}'

def _has_img_xml(xml):
    """Detecta qualquer tipo de imagem (DrawingML moderno ou VML legado)."""
    return ('<w:drawing' in xml) or ('<w:pict' in xml) or ('<v:imagedata' in xml)

def _img_asset(k):
    p = {"capa_background": "capa_background.png",
         "logo_carranza": "logo_carranza.png",
         "encerramento": "encerramento.png"}[k]
    with open(os.path.join(ASSETS, p), "rb") as f:
        return io.BytesIO(f.read())

def _blank(prs):
    for lay in prs.slide_layouts:
        if lay.name.lower() in ("em branco","blank",""): return lay
    return prs.slide_layouts[6]

def _pic(slide, k, l, t, w, h):
    b = _img_asset(k); b.seek(0)
    return slide.shapes.add_picture(b, Emu(l), Emu(t), Emu(w), Emu(h))

def _run(para, txt, bold=False, sz=635000, color=None):
    r = para.add_run(); r.text = txt
    r.font.bold = bold; r.font.size = Emu(sz); r.font.name = "Calibri"
    if color: r.font.color.rgb = color
    return r

def _faixa_rodape(slide):
    """Adiciona faixa vinho no rodapé de slides de conteúdo."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(0), Emu(FAIXA_Y), Emu(SLIDE_W), Emu(FAIXA_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x70, 0x00, 0x1C)
    shape.line.fill.background()  # sem borda

def _slide_capa(prs, dados):
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"capa_background",0,0,SLIDE_W,SLIDE_H)
    _pic(s,"logo_carranza",LOGO_CAPA_X,LOGO_CAPA_Y,LOGO_CAPA_CX,LOGO_CAPA_CY)
    tb = s.shapes.add_textbox(Emu(CAPA_X),Emu(CAPA_Y),Emu(CAPA_CX),Emu(CAPA_CY))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for txt, bold, sz in [
        (dados.get("disciplina","").upper(), True, 635000),
        (dados.get("assunto","").upper(), True, 635000),
        (dados.get("tipo","QUESTOES").upper(), True, 508000),
        (dados.get("professor",""), False, 355600),
    ]:
        if not txt: continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.alignment = PP_ALIGN.CENTER
        _run(p, txt, bold=bold, sz=sz, color=VINHO)

def _set_space_before(para, pts):
    """Seta space-before em um a:p do pptx via XML (python-pptx nao tem atalho)."""
    try:
        from pptx.oxml.ns import qn as _qn
        from lxml import etree as _et
        pPr = para._p.get_or_add_pPr()
        # Remover spcBef existente
        for old in pPr.findall(_qn('a:spcBef')):
            pPr.remove(old)
        spcBef = _et.SubElement(pPr, _qn('a:spcBef'))
        spcPts = _et.SubElement(spcBef, _qn('a:spcPts'))
        # spcPts val em centipontos (1pt = 100)
        spcPts.set('val', str(int(pts * 100)))
        # Garantir que o spcBef venha ANTES de elementos de run como buChar etc.
        # Mover pra primeira posicao do pPr eh mais seguro.
        pPr.insert(0, spcBef)
    except Exception:
        pass

def _slide_conteudo(prs, paragrafos, citacao=None):
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"logo_carranza",LOGO_MASTER_X,LOGO_MASTER_Y,LOGO_MASTER_CX,LOGO_MASTER_CY)
    tb = s.shapes.add_textbox(Emu(TEXTO_X),Emu(TEXTO_Y),Emu(TEXTO_CX),Emu(TEXTO_H))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for p in paragrafos:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; para.alignment = PP_ALIGN.JUSTIFY
        # Espaco antes (opcional) - separa enunciado da primeira alternativa
        sb_pt = p.get("space_before_pt")
        if sb_pt:
            _set_space_before(para, sb_pt)
        _run(para, p.get("text",""), bold=p.get("bold",False),
             sz=p.get("sz",571500), color=p.get("color",PRETO))
    if citacao:
        tb2 = s.shapes.add_textbox(Emu(RODAPE_X),Emu(RODAPE_Y),Emu(RODAPE_CX),Emu(RODAPE_CY))
        p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
        _run(p2, citacao, bold=True, sz=254000, color=VINHO)
    _faixa_rodape(s)

def _slide_imagem(prs, img_bytes_b64, img_ext="png"):
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"logo_carranza",LOGO_MASTER_X,LOGO_MASTER_Y,LOGO_MASTER_CX,LOGO_MASTER_CY)
    # Area util abaixo da logo master (evita sobreposicao)
    margin_x = TEXTO_X
    margin_y = LOGO_MASTER_Y + LOGO_MASTER_CY + 120000
    max_w = SLIDE_W - 2 * margin_x
    max_h = RODAPE_Y - margin_y - int(20 * 12700)
    img_bytes = base64.b64decode(img_bytes_b64)
    # Fit-to-area mantendo aspect ratio — SEM cap em 1.0, ou seja,
    # amplia a imagem ate ocupar o maximo do slide sem distorcer.
    w_emu = int(max_w * 0.9); h_emu = int(max_h * 0.9)  # fallback
    try:
        from PIL import Image as PILImage
        pil = PILImage.open(io.BytesIO(img_bytes))
        orig_w, orig_h = pil.size
        if orig_w > 0 and orig_h > 0:
            ratio = orig_w / orig_h
            # Descobrir qual dimensao limita (largura ou altura)
            if max_w / max_h >= ratio:
                # Altura eh o limite -> ocupa max_h e calcula largura proporcional
                h_emu = max_h
                w_emu = int(max_h * ratio)
            else:
                # Largura eh o limite -> ocupa max_w e calcula altura proporcional
                w_emu = max_w
                h_emu = int(max_w / ratio)
    except Exception:
        pass
    left = (SLIDE_W - w_emu) // 2
    top = margin_y + (max_h - h_emu) // 2
    buf = io.BytesIO(img_bytes); buf.seek(0)
    s.shapes.add_picture(buf, Emu(left), Emu(top), Emu(w_emu), Emu(h_emu))
    _faixa_rodape(s)

def _slide_gabarito(prs, qs, rs):
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"logo_carranza",LOGO_MASTER_X,LOGO_MASTER_Y,LOGO_MASTER_CX,LOGO_MASTER_CY)
    tb = s.shapes.add_textbox(Emu(4000000),Emu(2400000),Emu(4200000),Emu(600000))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "GABARITO", bold=True, sz=304800, color=VINHO)
    n = len(qs)
    if not n: return
    CINZA_ESCURO = RGBColor(0x59, 0x59, 0x59)  # Office Gray 60%
    CINZA_CLARO  = RGBColor(0xD9, 0xD9, 0xD9)  # Office Gray 15%
    BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
    tbl = s.shapes.add_table(2, n, Emu(2031997), Emu(3058160), Emu(8128005), Emu(741680)).table
    for j, (q, r) in enumerate(zip(qs, rs)):
        # Linha superior: numero da questao em fundo cinza escuro, texto branco
        c0 = tbl.cell(0,j); c0.text = str(q).zfill(2)
        try:
            c0.fill.solid(); c0.fill.fore_color.rgb = CINZA_ESCURO
        except Exception: pass
        p0 = c0.text_frame.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        for ru in p0.runs:
            ru.font.bold = True; ru.font.size = Pt(18); ru.font.color.rgb = BRANCO
        # Linha inferior: resposta em fundo cinza claro, texto preto
        c1 = tbl.cell(1,j); c1.text = str(r).upper()
        try:
            c1.fill.solid(); c1.fill.fore_color.rgb = CINZA_CLARO
        except Exception: pass
        p1 = c1.text_frame.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        for ru in p1.runs:
            ru.font.bold = True; ru.font.size = Pt(18); ru.font.color.rgb = PRETO
    _faixa_rodape(s)

def _slide_secao(prs, titulo):
    """Slide divisor de seção — fundo da capa com texto centralizado."""
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"capa_background",0,0,SLIDE_W,SLIDE_H)
    _pic(s,"logo_carranza",LOGO_CAPA_X,LOGO_CAPA_Y,LOGO_CAPA_CX,LOGO_CAPA_CY)
    tb = s.shapes.add_textbox(Emu(CAPA_X),Emu(CAPA_Y),Emu(CAPA_CX),Emu(CAPA_CY))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, titulo.upper(), bold=True, sz=635000, color=VINHO)

def _slide_conteudo_titulado(prs, titulo, paragrafos, citacao=None):
    """Slide de conteúdo com título em destaque + parágrafos abaixo."""
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"logo_carranza",LOGO_MASTER_X,LOGO_MASTER_Y,LOGO_MASTER_CX,LOGO_MASTER_CY)
    # Título: abaixo da faixa, NÃO sobrepõe a logo (largura limitada)
    TIT_Y = 350000
    TIT_W = LOGO_MASTER_X - TEXTO_X - 200000   # para antes da logo
    TIT_H = 650000
    tb_titulo = s.shapes.add_textbox(Emu(TEXTO_X), Emu(TIT_Y), Emu(TIT_W), Emu(TIT_H))
    tf_t = tb_titulo.text_frame; tf_t.word_wrap = True
    pt = tf_t.paragraphs[0]; pt.alignment = PP_ALIGN.LEFT
    _run(pt, titulo.upper() if titulo else "", bold=True, sz=508000, color=VINHO)
    # Conteúdo: começa ABAIXO da logo
    CONT_Y = LOGO_MASTER_Y + LOGO_MASTER_CY + 80000  # logo bottom + margem
    CONT_H = RODAPE_Y - CONT_Y - int(20 * 12700)
    tb = s.shapes.add_textbox(Emu(TEXTO_X), Emu(CONT_Y), Emu(TEXTO_CX), Emu(CONT_H))
    tf = tb.text_frame; tf.word_wrap = True
    # Calcular tamanho da fonte baseado no volume de texto
    total_chars = sum(len(p) for p in paragrafos)
    if total_chars > 800: sz = 381000
    elif total_chars > 500: sz = 406400
    elif total_chars > 300: sz = 444500
    else: sz = 482600
    first = True
    for texto in paragrafos:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; para.alignment = PP_ALIGN.JUSTIFY
        _run(para, texto, bold=False, sz=sz, color=PRETO)
    if citacao:
        tb2 = s.shapes.add_textbox(Emu(RODAPE_X),Emu(RODAPE_Y),Emu(RODAPE_CX),Emu(RODAPE_CY))
        p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
        _run(p2, citacao, bold=True, sz=254000, color=VINHO)
    _faixa_rodape(s)

def _slide_tabela(prs, n_rows, n_cols, cells, titulo=None):
    """Cria um slide com uma tabela no padrao Carranza.

    cells = [{"r": int, "c": int, "rowspan": int, "colspan": int, "text": str}, ...]
    Cada item representa uma celula MASTER (os merges ja foram resolvidos).
    """
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"logo_carranza",LOGO_MASTER_X,LOGO_MASTER_Y,LOGO_MASTER_CX,LOGO_MASTER_CY)
    if n_rows <= 0 or n_cols <= 0 or not cells:
        _faixa_rodape(s); return

    # Area util abaixo da logo
    margin_x = TEXTO_X
    margin_y = LOGO_MASTER_Y + LOGO_MASTER_CY + 200000
    # Titulo opcional
    if titulo:
        tbx = s.shapes.add_textbox(Emu(margin_x), Emu(margin_y), Emu(SLIDE_W - 2*margin_x), Emu(450000))
        p_t = tbx.text_frame.paragraphs[0]; p_t.alignment = PP_ALIGN.CENTER
        _run(p_t, titulo.upper(), bold=True, sz=355600, color=VINHO)
        margin_y += 500000
    max_w = SLIDE_W - 2 * margin_x
    max_h = RODAPE_Y - margin_y - int(30 * 12700)

    # Altura proporcional ao numero de linhas
    row_h = max(400000, min(800000, max_h // max(1, n_rows)))
    tbl_h = min(max_h, row_h * n_rows)

    table_shape = s.shapes.add_table(n_rows, n_cols,
                                     Emu(margin_x), Emu(margin_y),
                                     Emu(max_w), Emu(tbl_h))
    tbl = table_shape.table

    # Fonte adaptativa pelo volume de texto
    total_chars = sum(len(c["text"]) for c in cells)
    if total_chars > 1500: font_sz = Pt(9)
    elif total_chars > 1000: font_sz = Pt(10)
    elif total_chars > 600:  font_sz = Pt(11)
    elif total_chars > 300:  font_sz = Pt(12)
    else:                     font_sz = Pt(14)

    # 1) Aplicar merges primeiro (pythonpptx: cell.merge(other))
    for c in cells:
        if c["rowspan"] > 1 or c["colspan"] > 1:
            r0, c0 = c["r"], c["c"]
            r1 = r0 + c["rowspan"] - 1
            c1 = c0 + c["colspan"] - 1
            try:
                tbl.cell(r0, c0).merge(tbl.cell(r1, c1))
            except Exception:
                pass

    # 2) Preencher conteudo e estilo em cada master
    for c in cells:
        r0, c0 = c["r"], c["c"]
        cell = tbl.cell(r0, c0)
        is_header = (r0 == 0)
        # Reset
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        # Margens internas menores para aproveitar espaco
        try:
            from pptx.util import Inches as _In
            cell.margin_left = _In(0.06); cell.margin_right = _In(0.06)
            cell.margin_top  = _In(0.03); cell.margin_bottom = _In(0.03)
        except Exception: pass
        lines = [ln for ln in c["text"].split("\n")] or [""]
        # Remove linhas vazias redundantes no topo
        while lines and not lines[0].strip(): lines.pop(0)
        while lines and not lines[-1].strip(): lines.pop()
        if not lines: lines = [""]
        first = True
        for ln in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT
            ru = p.add_run(); ru.text = ln
            ru.font.size = font_sz
            ru.font.name = "Calibri"
            if is_header:
                ru.font.bold = True
                ru.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                ru.font.color.rgb = PRETO
        # Fundo do header em vinho
        if is_header:
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = VINHO
            except Exception: pass

    _faixa_rodape(s)

def _slide_enc(prs):
    s = prs.slides.add_slide(_blank(prs))
    _pic(s,"capa_background",0,0,SLIDE_W,SLIDE_H)
    _pic(s,"logo_carranza",LOGO_ENC_X,LOGO_ENC_Y,LOGO_ENC_CX,LOGO_ENC_CY)

def _h(txt, sz):
    """Estima altura do texto em EMU, respeitando quebras de linha (\n)."""
    if not txt.strip(): return int(sz*0.4)
    cpp = LARGURA / (sz * FATOR)
    linhas = 0
    for seg in txt.split("\n"):
        linhas += max(1, int((len(seg) + cpp - 1) // cpp))
    # \n\n adiciona um respiro extra entre paragrafos (~0.4 linha)
    linhas += txt.count("\n\n") * 0.4
    return int(sz * 1.15 * max(1.0, linhas))

def _sz(te, alts):
    total = len(te) + sum(len(a) for a in alts)
    if total > 1500: return 330200   # 26pt  (enunciados muito longos)
    if total > 1000: return 355600   # 28pt
    if total > 800:  return 406400   # 32pt
    if total > 500:  return 444500   # 35pt
    if total > 300:  return 482600   # 38pt
    return 533400                    # 42pt

def _split_enunciado(te, sz):
    """Divide o enunciado em chunks que cabem em AREA_UTIL.

    Quebra primeiro por paragrafos (linha em branco); se um paragrafo ainda
    for grande demais, quebra por frases; em ultimo caso faz corte bruto."""
    import re as _re
    if _h(te, sz) <= AREA_UTIL:
        return [te]
    paras = [p.strip() for p in _re.split(r"\n\s*\n", te) if p.strip()]
    if not paras:
        paras = [te]

    def _empacotar_frases(p):
        frases = _re.split(r"(?<=[\.\!\?])\s+", p)
        buf, out = "", []
        for f in frases:
            tent = (buf + " " + f) if buf else f
            if _h(tent, sz) <= AREA_UTIL:
                buf = tent
            else:
                if buf: out.append(buf)
                if _h(f, sz) > AREA_UTIL:
                    # Corte bruto proporcional ao excesso
                    chars_max = max(1, int(len(f) * AREA_UTIL / _h(f, sz)))
                    for i in range(0, len(f), chars_max):
                        out.append(f[i:i+chars_max])
                    buf = ""
                else:
                    buf = f
        if buf: out.append(buf)
        return out

    chunks, atual = [], ""
    for p in paras:
        tentativa = (atual + "\n\n" + p) if atual else p
        if _h(tentativa, sz) <= AREA_UTIL:
            atual = tentativa
        else:
            if atual:
                chunks.append(atual); atual = ""
            if _h(p, sz) > AREA_UTIL:
                chunks.extend(_empacotar_frases(p))
            else:
                atual = p
    if atual: chunks.append(atual)
    return chunks


def _distribuir(te, alts, sz):
    """Distribui enunciado + alternativas em grupos (= slides) que cabem.

    Correcao: o enunciado agora eh quebrado em multiplos slides se nao
    couber em um so. Antes, um enunciado gigante era empurrado inteiro
    num unico slide e estourava a margem inferior.
    """
    grupos = []
    enunciado_chunks = _split_enunciado(te, sz)
    # Chunks iniciais do enunciado viram slides dedicados
    for ch in enunciado_chunks[:-1]:
        grupos.append([{"text": ch, "bold": True, "sz": sz}])
    # Ultimo chunk do enunciado + alternativas (com quebra por altura)
    grupo = [{"text": enunciado_chunks[-1], "bold": True, "sz": sz}]
    altura = _h(enunciado_chunks[-1], sz)
    primeira_alt = True
    for alt in alts:
        ha = _h(alt, sz)
        if primeira_alt:
            entry = {"text": alt, "bold": False, "sz": sz, "space_before_pt": 14}
            extra = 14 * 12700
        else:
            entry = {"text": alt, "bold": False, "sz": sz}
            extra = 0
        if altura + ha + extra > AREA_UTIL and grupo:
            grupos.append(grupo); grupo = []; altura = 0
            entry.pop("space_before_pt", None)
        grupo.append(entry)
        altura += ha + extra
        primeira_alt = False
    if grupo: grupos.append(grupo)
    return grupos

def _gerar_docx(payload):
    """
    Gera um .docx formatado no padrao Carranza a partir de conteudo extraido de slides.
    Usa o MODELO DOCUMENTO.docx como base, preservando:
    - Capa com imagem de fundo (aguia)
    - Header com logo Carranza Concursos
    - Footer com barra de redes sociais
    Retorna um BytesIO com o .docx pronto.
    """
    import copy as copymod

    # Copiar template para temp
    tmp_docx = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
    tmp_docx.close()
    shutil.copy2(MODELO_DOCX, tmp_docx.name)

    doc = DocxDocument(tmp_docx.name)
    body = doc.element.body
    children = list(body)
    sect_pr = body.find(qn('w:sectPr'))

    # ---- Identificar elementos da capa ----
    # [0] = Background image text box (full-page aguia)
    # [9] = Title text box (onde vai disciplina/assunto)
    # Salvar capa (paragrafos 0 a 9) e remover o resto
    WPS_TXBX = '{http://schemas.microsoft.com/office/word/2010/wordprocessingShape}txbx'
    V_TEXTBOX = '{urn:schemas-microsoft-com:vml}textbox'

    cover_elements = []  # paragrafos da capa para preservar
    title_textbox_para = None  # paragrafo com text box do titulo
    bg_para = None  # paragrafo com background

    drawing_count = 0
    for i, child in enumerate(children):
        if child.tag == qn('w:sectPr'):
            continue
        has_drawing = len(child.findall('.//' + qn('w:drawing'))) > 0
        if has_drawing:
            drawing_count += 1
            if drawing_count == 1:
                bg_para = child  # primeiro = background
                cover_elements.append(child)
            elif drawing_count == 2:
                title_textbox_para = child  # segundo = text box titulo
                cover_elements.append(child)
            # Ignorar drawings 3+ (sobras do template)
        else:
            # Paragrafos vazios entre background e titulo fazem parte da capa
            if drawing_count >= 1 and drawing_count <= 2 and i < 15:
                cover_elements.append(child)

    # (bg_para_copy removido - pagina final usa imagem diretamente)

    # ---- Injetar texto no text box do titulo ----
    disc = payload.get("disciplina", "DISCIPLINA")
    ass = payload.get("assunto", "")
    prof = payload.get("professor", "")
    tipo = payload.get("tipo", "QUESTOES")

    if title_textbox_para is not None:
        # Encontrar o txbxContent dentro do text box
        txbx_contents = title_textbox_para.findall('.//' + WPS_TXBX)
        if not txbx_contents:
            txbx_contents = title_textbox_para.findall('.//' + V_TEXTBOX)

        for txbx in txbx_contents:
            # Encontrar w:txbxContent
            for txbc in txbx.findall('.//' + qn('w:txbxContent')):
                # Limpar paragrafos existentes
                for old_p in list(txbc):
                    txbc.remove(old_p)

                # Adicionar titulo da disciplina
                lines = []
                if disc: lines.append((disc.upper(), "56", True))
                if ass: lines.append((ass.upper(), "40", True))
                if tipo: lines.append((tipo.upper(), "28", False))
                if prof: lines.append((prof, "24", False))

                for text, sz, bold in lines:
                    p_el = etree.SubElement(txbc, qn('w:p'))
                    pPr = etree.SubElement(p_el, qn('w:pPr'))
                    sp = etree.SubElement(pPr, qn('w:spacing'))
                    sp.set(qn('w:after'), '80')
                    sp.set(qn('w:line'), '240')
                    sp.set(qn('w:lineRule'), 'auto')
                    jc = etree.SubElement(pPr, qn('w:jc'))
                    jc.set(qn('w:val'), 'center')
                    r_el = etree.SubElement(p_el, qn('w:r'))
                    rPr = etree.SubElement(r_el, qn('w:rPr'))
                    if bold:
                        etree.SubElement(rPr, qn('w:b'))
                        etree.SubElement(rPr, qn('w:bCs'))
                    color = etree.SubElement(rPr, qn('w:color'))
                    color.set(qn('w:val'), '70001C')
                    sz_el = etree.SubElement(rPr, qn('w:sz'))
                    sz_el.set(qn('w:val'), sz)
                    szCs = etree.SubElement(rPr, qn('w:szCs'))
                    szCs.set(qn('w:val'), sz)
                    fn = etree.SubElement(rPr, qn('w:rFonts'))
                    fn.set(qn('w:ascii'), 'Calibri')
                    fn.set(qn('w:hAnsi'), 'Calibri')
                    t_el = etree.SubElement(r_el, qn('w:t'))
                    t_el.text = text
                break  # processar apenas o primeiro txbx

    # ---- Remover TUDO do body exceto capa + sectPr ----
    for child in list(body):
        if child.tag == qn('w:sectPr'):
            continue
        body.remove(child)

    # Re-inserir elementos da capa antes do sectPr
    for elem in cover_elements:
        body.insert(list(body).index(sect_pr), elem)

    # ---- Constantes de estilo ----
    VINHO_HEX = "70001C"
    FONT_NAME = "Calibri"

    def _add_para_before_sectpr(text="", font_size=11, bold=False, color_hex=None,
                                alignment=None, space_before=0, space_after=120,
                                font_name=FONT_NAME, italic=False):
        """Adiciona um paragrafo estilizado ANTES do sectPr."""
        p_el = etree.SubElement(body, qn('w:p'))
        # Mover para antes do sectPr
        body.remove(p_el)
        body.insert(list(body).index(sect_pr), p_el)

        # paragraph properties
        pPr = etree.SubElement(p_el, qn('w:pPr'))
        sp = etree.SubElement(pPr, qn('w:spacing'))
        sp.set(qn('w:before'), str(int(space_before * 20)))
        sp.set(qn('w:after'), str(int(space_after * 20)))
        if alignment is not None:
            jc = etree.SubElement(pPr, qn('w:jc'))
            align_map = {WD_ALIGN_PARAGRAPH.CENTER: 'center', WD_ALIGN_PARAGRAPH.LEFT: 'left',
                         WD_ALIGN_PARAGRAPH.RIGHT: 'right', WD_ALIGN_PARAGRAPH.JUSTIFY: 'both'}
            jc.set(qn('w:val'), align_map.get(alignment, 'left'))

        if not text:
            return p_el

        r_el = etree.SubElement(p_el, qn('w:r'))
        rPr = etree.SubElement(r_el, qn('w:rPr'))
        if bold:
            etree.SubElement(rPr, qn('w:b'))
            etree.SubElement(rPr, qn('w:bCs'))
        if italic:
            etree.SubElement(rPr, qn('w:i'))
            etree.SubElement(rPr, qn('w:iCs'))
        if color_hex:
            c = etree.SubElement(rPr, qn('w:color'))
            c.set(qn('w:val'), color_hex)
        sz_val = str(int(font_size * 2))
        sz_el = etree.SubElement(rPr, qn('w:sz'))
        sz_el.set(qn('w:val'), sz_val)
        szCs = etree.SubElement(rPr, qn('w:szCs'))
        szCs.set(qn('w:val'), sz_val)
        fn = etree.SubElement(rPr, qn('w:rFonts'))
        fn.set(qn('w:ascii'), font_name)
        fn.set(qn('w:hAnsi'), font_name)
        t_el = etree.SubElement(r_el, qn('w:t'))
        t_el.text = text
        t_el.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')

        return p_el

    def _add_pagebreak():
        """Adiciona quebra de pagina antes do sectPr."""
        p_el = etree.SubElement(body, qn('w:p'))
        body.remove(p_el)
        body.insert(list(body).index(sect_pr), p_el)
        r_el = etree.SubElement(p_el, qn('w:r'))
        br = etree.SubElement(r_el, qn('w:br'))
        br.set(qn('w:type'), 'page')
        return p_el

    def _add_separator():
        """Adiciona linha horizontal vinho."""
        p_el = _add_para_before_sectpr("", space_before=2, space_after=2)
        pPr = p_el.find(qn('w:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p_el, qn('w:pPr'))
        pBdr = etree.SubElement(pPr, qn('w:pBdr'))
        bottom = etree.SubElement(pBdr, qn('w:bottom'))
        bottom.set(qn('w:val'), 'single')
        bottom.set(qn('w:sz'), '12')
        bottom.set(qn('w:space'), '1')
        bottom.set(qn('w:color'), VINHO_HEX)
        return p_el

    def _add_questao_para(numero, enunciado):
        """Adiciona paragrafo de questao com numero em vinho + enunciado."""
        p_el = etree.SubElement(body, qn('w:p'))
        body.remove(p_el)
        body.insert(list(body).index(sect_pr), p_el)

        pPr = etree.SubElement(p_el, qn('w:pPr'))
        sp = etree.SubElement(pPr, qn('w:spacing'))
        sp.set(qn('w:before'), '200')
        sp.set(qn('w:after'), '80')

        # Run do numero (vinho, bold)
        r1 = etree.SubElement(p_el, qn('w:r'))
        rPr1 = etree.SubElement(r1, qn('w:rPr'))
        etree.SubElement(rPr1, qn('w:b'))
        etree.SubElement(rPr1, qn('w:bCs'))
        c1 = etree.SubElement(rPr1, qn('w:color'))
        c1.set(qn('w:val'), VINHO_HEX)
        sz1 = etree.SubElement(rPr1, qn('w:sz'))
        sz1.set(qn('w:val'), '22')
        szCs1 = etree.SubElement(rPr1, qn('w:szCs'))
        szCs1.set(qn('w:val'), '22')
        fn1 = etree.SubElement(rPr1, qn('w:rFonts'))
        fn1.set(qn('w:ascii'), 'Calibri')
        fn1.set(qn('w:hAnsi'), 'Calibri')
        t1 = etree.SubElement(r1, qn('w:t'))
        t1.text = f"{numero}. "
        t1.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')

        # Run do enunciado (preto)
        r2 = etree.SubElement(p_el, qn('w:r'))
        rPr2 = etree.SubElement(r2, qn('w:rPr'))
        c2 = etree.SubElement(rPr2, qn('w:color'))
        c2.set(qn('w:val'), '333333')
        sz2 = etree.SubElement(rPr2, qn('w:sz'))
        sz2.set(qn('w:val'), '22')
        szCs2 = etree.SubElement(rPr2, qn('w:szCs'))
        szCs2.set(qn('w:val'), '22')
        fn2 = etree.SubElement(rPr2, qn('w:rFonts'))
        fn2.set(qn('w:ascii'), 'Calibri')
        fn2.set(qn('w:hAnsi'), 'Calibri')
        t2 = etree.SubElement(r2, qn('w:t'))
        t2.text = enunciado
        return p_el

    # ---- Helpers para tipos adicionais (tabela, imagem) ----
    def _sombrear_celula(cell, hex_fill):
        """Aplica fundo colorido em uma celula (w:shd)."""
        tcPr = cell._tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('w:shd')):
            tcPr.remove(old)
        shd = etree.SubElement(tcPr, qn('w:shd'))
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_fill)

    def _aplicar_bordas_tabela(tbl, color='BFBFBF', sz='4'):
        """Bordas finas cinzas em todos os lados + entre celulas."""
        tblPr = tbl._tbl.find(qn('w:tblPr'))
        if tblPr is None:
            tblPr = etree.SubElement(tbl._tbl, qn('w:tblPr'))
            tbl._tbl.insert(0, tblPr)
        borders = tblPr.find(qn('w:tblBorders'))
        if borders is None:
            borders = etree.SubElement(tblPr, qn('w:tblBorders'))
        for b_name in ('top','left','bottom','right','insideH','insideV'):
            tag = qn('w:' + b_name)
            b = borders.find(tag)
            if b is None:
                b = etree.SubElement(borders, tag)
            b.set(qn('w:val'), 'single')
            b.set(qn('w:sz'), sz)
            b.set(qn('w:color'), color)

    def _estilizar_celula(cell, text, is_header):
        """Preenche celula com texto multi-linha e aplica tipografia."""
        lines = [ln for ln in (text or "").split("\n")]
        # Remove linhas vazias no topo/fim
        while lines and not lines[0].strip(): lines.pop(0)
        while lines and not lines[-1].strip(): lines.pop()
        if not lines: lines = [""]
        # Resetar o primeiro paragrafo (que sempre existe numa celula recem-criada)
        cell.text = ""
        first_p = cell.paragraphs[0]
        def _fmt(p, text):
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER if is_header else WD_ALIGN_PARAGRAPH.LEFT
            run = p.add_run(text)
            run.font.name = 'Calibri'
            run.font.size = DocxPt(10 if is_header else 10)
            if is_header:
                run.font.bold = True
                run.font.color.rgb = DocxRGB(0xFF, 0xFF, 0xFF)
            else:
                run.font.color.rgb = DocxRGB(0x33, 0x33, 0x33)
        _fmt(first_p, lines[0])
        for ln in lines[1:]:
            _fmt(cell.add_paragraph(), ln)

    def _add_tabela_docx(n_rows, n_cols, cells, titulo=None):
        """Cria uma tabela com merges e cabecalho vinho, movendo para antes do sectPr."""
        if n_rows <= 0 or n_cols <= 0 or not cells:
            return None
        # Titulo opcional antes da tabela
        if titulo:
            _add_para_before_sectpr(titulo.upper(), font_size=12, bold=True,
                                    color_hex=VINHO_HEX, alignment=WD_ALIGN_PARAGRAPH.CENTER,
                                    space_before=8, space_after=4)
        tbl = doc.add_table(rows=n_rows, cols=n_cols)
        # Larga 100%
        try:
            from docx.shared import Inches as _DocxInches
            tbl.autofit = False
            tblPr = tbl._tbl.find(qn('w:tblPr'))
            if tblPr is None:
                tblPr = etree.SubElement(tbl._tbl, qn('w:tblPr'))
                tbl._tbl.insert(0, tblPr)
            tblW = tblPr.find(qn('w:tblW'))
            if tblW is None:
                tblW = etree.SubElement(tblPr, qn('w:tblW'))
            tblW.set(qn('w:w'), '5000')
            tblW.set(qn('w:type'), 'pct')
        except Exception:
            pass
        _aplicar_bordas_tabela(tbl)
        # 1) Aplicar merges primeiro
        for c in cells:
            if c.get("rowspan", 1) > 1 or c.get("colspan", 1) > 1:
                r0, c0 = c["r"], c["c"]
                r1 = r0 + c["rowspan"] - 1
                c1 = c0 + c["colspan"] - 1
                try: tbl.cell(r0, c0).merge(tbl.cell(r1, c1))
                except Exception: pass
        # 2) Preencher conteudo
        for c in cells:
            r0, c0 = c["r"], c["c"]
            cell = tbl.cell(r0, c0)
            is_header = (r0 == 0)
            if is_header:
                _sombrear_celula(cell, VINHO_HEX)
            _estilizar_celula(cell, c.get("text", ""), is_header)
        # 3) Mover para antes do sectPr
        body.remove(tbl._element)
        body.insert(list(body).index(sect_pr), tbl._element)
        # Espaco depois da tabela
        _add_para_before_sectpr("", font_size=4, space_after=6)
        return tbl

    def _add_imagem_docx(img_b64, img_ext="png"):
        """Insere imagem centralizada, fit-to-area A4 mantendo ratio."""
        if not img_b64: return None
        try:
            img_bytes = base64.b64decode(img_b64)
        except Exception:
            return None
        # Area util A4 retrato com margens padrao: ~16cm x 22cm
        max_w_cm = 15.5
        max_h_cm = 20.0
        w_cm, h_cm = max_w_cm * 0.7, max_h_cm * 0.5  # fallback
        try:
            from PIL import Image as PILImage
            pil = PILImage.open(io.BytesIO(img_bytes))
            ow, oh = pil.size
            if ow > 0 and oh > 0:
                ratio = ow / oh
                if max_w_cm / max_h_cm >= ratio:
                    h_cm = max_h_cm
                    w_cm = max_h_cm * ratio
                else:
                    w_cm = max_w_cm
                    h_cm = max_w_cm / ratio
        except Exception:
            pass
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        try:
            run.add_picture(io.BytesIO(img_bytes), width=Cm(w_cm), height=Cm(h_cm))
        except Exception:
            return None
        body.remove(p._element)
        body.insert(list(body).index(sect_pr), p._element)
        _add_para_before_sectpr("", font_size=4, space_after=4)
        return p

    # ---- Quebra de pagina apos capa ----
    _add_pagebreak()

    # ---- Conteudo dos slides ----
    slides = payload.get("slides", [])
    gabarito = payload.get("gabarito")
    num_questao = 0

    for sl in slides:
        t = sl.get("tipo", "")

        if t == "secao":
            _add_separator()
            _add_para_before_sectpr(sl.get("titulo", "").upper(), font_size=14, bold=True,
                                    color_hex=VINHO_HEX, space_before=6, space_after=6,
                                    alignment=WD_ALIGN_PARAGRAPH.LEFT)
            _add_separator()
            _add_para_before_sectpr("", font_size=4, space_after=4)

        elif t == "conteudo_slide":
            titulo = sl.get("titulo", "")
            if titulo:
                _add_para_before_sectpr(titulo, font_size=12, bold=True,
                                        color_hex=VINHO_HEX, space_before=10, space_after=4)
            for para_text in sl.get("paragrafos", []):
                if para_text.strip():
                    _add_para_before_sectpr(para_text, font_size=11, bold=False,
                                            color_hex="333333", space_before=0, space_after=3,
                                            alignment=WD_ALIGN_PARAGRAPH.JUSTIFY)

        elif t == "questao":
            num_questao += 1
            numero = sl.get("numero", num_questao)
            enunciado = sl.get("enunciado", "")
            _add_questao_para(numero, enunciado)

            if sl.get("certo_errado", False):
                _add_para_before_sectpr("(   ) Certo    (   ) Errado", font_size=11,
                                        color_hex="333333", space_before=1, space_after=4)
            else:
                for alt in sl.get("alternativas", []):
                    _add_para_before_sectpr(alt, font_size=11, color_hex="333333",
                                            space_before=0, space_after=1)

        elif t == "contexto":
            texto = sl.get("texto", "")
            if texto:
                _add_para_before_sectpr(texto, font_size=11, color_hex="333333",
                                        space_before=4, space_after=4,
                                        alignment=WD_ALIGN_PARAGRAPH.JUSTIFY)

        elif t == "tabela":
            _add_tabela_docx(
                sl.get("n_rows", 0),
                sl.get("n_cols", 0),
                sl.get("cells", []),
                titulo=sl.get("titulo"),
            )

        elif t == "imagem":
            _add_imagem_docx(sl.get("img_b64", ""), sl.get("img_ext", "png"))

    # ---- Gabarito ----
    if gabarito and gabarito.get("questoes"):
        _add_pagebreak()
        _add_separator()
        _add_para_before_sectpr("GABARITO", font_size=16, bold=True,
                                color_hex=VINHO_HEX, alignment=WD_ALIGN_PARAGRAPH.CENTER,
                                space_before=6, space_after=6)
        _add_separator()
        _add_para_before_sectpr("", font_size=4, space_after=4)
        qs = gabarito.get("questoes", [])
        rs = gabarito.get("respostas", [])
        for i in range(0, len(qs), 5):
            chunk_q = qs[i:i+5]
            chunk_r = rs[i:i+5]
            parts = [f"{q}) {r}" for q, r in zip(chunk_q, chunk_r)]
            _add_para_before_sectpr("     ".join(parts), font_size=12, bold=True,
                                    color_hex=VINHO_HEX, alignment=WD_ALIGN_PARAGRAPH.CENTER,
                                    space_before=0, space_after=3)

    # ---- Pagina final (encerramento) com imagem de fundo ----
    _add_pagebreak()
    try:
        bg_path = os.path.join(ASSETS, "doc_background.jpeg")
        if os.path.exists(bg_path):
            # Adicionar imagem de fundo centralizada como pagina de encerramento
            p_img = doc.add_paragraph()
            p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run_img = p_img.add_run()
            # A4 ajustado para nao sobrepor header/footer
            run_img.add_picture(bg_path, width=Cm(14.5), height=Cm(20.5))
            pf = p_img.paragraph_format
            pf.space_before = DocxPt(0)
            pf.space_after = DocxPt(0)
            # Mover para antes do sectPr
            body.remove(p_img._element)
            body.insert(list(body).index(sect_pr), p_img._element)
    except Exception:
        pass  # se falhar, nao adiciona pagina final

    # ---- Salvar ----
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    try: os.unlink(tmp_docx.name)
    except: pass

    return buf


def _build_pptx(payload):
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W); prs.slide_height = Emu(SLIDE_H)
    dados = {k: payload.get(k,"") for k in ("disciplina","assunto","tipo","professor")}
    if not dados["tipo"]: dados["tipo"] = "QUESTOES"
    _slide_capa(prs, dados)
    cit = cycle(CITACOES)
    for s in payload.get("slides", []):
        tipo = s.get("tipo")
        if tipo == "contexto":
            _slide_conteudo(prs, [{"text": s.get("texto",""), "bold": True, "sz": 508000}])
            continue
        if tipo == "secao":
            _slide_secao(prs, s.get("titulo",""))
            continue
        if tipo == "conteudo_slide":
            _slide_conteudo_titulado(prs, s.get("titulo",""), s.get("paragrafos",[]), citacao=next(cit))
            continue
        if tipo == "tabela":
            _slide_tabela(prs,
                          s.get("n_rows", 0),
                          s.get("n_cols", 0),
                          s.get("cells", []),
                          titulo=s.get("titulo"))
            continue
        if tipo == "imagem":
            img_b64 = s.get("img_b64",""); img_ext = s.get("img_ext","png")
            if img_b64: _slide_imagem(prs, img_b64, img_ext)
            continue
        num = s.get("numero",""); enc = s.get("enunciado","")
        alts = s.get("alternativas",[]); ce = s.get("certo_errado", False)
        snum = str(num).zfill(2) if num else ""
        te = enc if (snum and (enc.startswith(snum+".") or enc.startswith(str(num)+"."))) else (snum+". "+enc if snum else enc)
        if ce:
            sz = _sz(te, ["Certo (  )","Errado (  )"])
            _slide_conteudo(prs, [
                {"text":te,"bold":True,"sz":sz},
                {"text":"Certo (  )","bold":False,"sz":sz,"space_before_pt":14},
                {"text":"Errado (  )","bold":False,"sz":sz},
            ], citacao=next(cit))
        elif alts:
            sz = _sz(te, alts)
            grupos = _distribuir(te, alts, sz)
            for gi, grupo in enumerate(grupos):
                _slide_conteudo(prs, grupo, citacao=next(cit) if gi==len(grupos)-1 else None)
        else:
            sz = _sz(te, [])
            _slide_conteudo(prs, [{"text":te,"bold":True,"sz":sz}], citacao=next(cit))
    gab = payload.get("gabarito")
    if gab and gab.get("questoes"):
        _slide_gabarito(prs, gab["questoes"], gab["respostas"])
    _slide_enc(prs)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf

def _rel_to_img(doc_part, rId):
    """Converte um relId em (b64, ext) ou None."""
    if not rId: return None
    try:
        rel = doc_part.rels.get(rId)
        if rel is None: return None
        ext = rel.target_ref.split('.')[-1].lower()
        if ext not in ('png','jpg','jpeg','gif','bmp','webp'): ext = 'png'
        return base64.b64encode(rel.target_part.blob).decode(), ext
    except Exception:
        return None

def _get_imgs_from_para(para, doc_part):
    """Extrai TODAS as imagens de um parágrafo (DrawingML + VML legado).

    Retorna lista de tuplas (b64, ext). Lista vazia se não houver imagens.
    """
    out = []
    xml = para._element.xml
    if not _has_img_xml(xml):
        return out
    seen = set()
    for blip in para._element.findall('.//' + NS_A + 'blip'):
        rId = blip.get(NS_R + 'embed') or blip.get(NS_R + 'link')
        if not rId or rId in seen: continue
        seen.add(rId)
        ir = _rel_to_img(doc_part, rId)
        if ir: out.append(ir)
    for imgdata in para._element.findall('.//' + VML_NS + 'imagedata'):
        rId = imgdata.get(NS_R + 'id') or imgdata.get(NS_R + 'embed')
        if not rId or rId in seen: continue
        seen.add(rId)
        ir = _rel_to_img(doc_part, rId)
        if ir: out.append(ir)
    return out

def _get_img_from_para(para, doc_part):
    """Wrapper de compatibilidade: retorna a primeira imagem ou None."""
    imgs = _get_imgs_from_para(para, doc_part)
    return imgs[0] if imgs else None

def _extrair_texto_docx(filepath):
    """Extrai texto bruto do docx para enviar ao Claude (compat)."""
    texto, _ = _extrair_texto_e_imgs_docx(filepath)
    return texto

def _tabela_eh_gabarito(tbl):
    """Detecta se uma tabela eh so gabarito (ex: celulas '01.E', '02.A').
    Retorna True se >= 70% das celulas nao-vazias batem padrao de gabarito.
    """
    RE_GAB = re.compile(r'^\d{1,2}\.\s*[A-Ea-eVvFfCcEe]\b')
    n_tot = 0; n_gab = 0
    for row in tbl.rows:
        for cell in row.cells:
            txt = cell.text.strip()
            if txt:
                n_tot += 1
                if RE_GAB.match(txt):
                    n_gab += 1
    return n_tot > 0 and (n_gab / n_tot) >= 0.7

def _extrair_tabela_como_dict(tbl):
    """Converte uma docx.Table em dict pronto para renderizar no PPTX.

    Resolve merges corretamente via OOXML (w:vMerge, w:gridSpan).
    Retorna {"n_rows", "n_cols", "cells": [{"r","c","rowspan","colspan","text"}]}.
    """
    tbl_el = tbl._tbl
    trs = tbl_el.findall(qn('w:tr'))
    n_rows = len(trs)
    if n_rows == 0:
        return None
    # Descobrir n_cols pelo tblGrid (quando disponivel) ou por soma de gridSpans
    tblGrid = tbl_el.find(qn('w:tblGrid'))
    if tblGrid is not None:
        n_cols = len(tblGrid.findall(qn('w:gridCol')))
    else:
        n_cols = 0
    if n_cols <= 0:
        n_cols = 0
        for tr in trs:
            col = 0
            for tc in tr.findall(qn('w:tc')):
                gs = 1
                tcPr = tc.find(qn('w:tcPr'))
                if tcPr is not None:
                    gsEl = tcPr.find(qn('w:gridSpan'))
                    if gsEl is not None:
                        try: gs = int(gsEl.get(qn('w:val'), '1'))
                        except Exception: pass
                col += gs
            if col > n_cols: n_cols = col
    if n_cols == 0:
        return None

    def _cell_text(tc):
        paras = tc.findall(qn('w:p'))
        lines = []
        for p in paras:
            ts = p.findall('.//' + qn('w:t'))
            lines.append(''.join(t.text or '' for t in ts))
        return '\n'.join(lines).strip()

    masters = []
    # Para cada coluna, qual master atualmente aberto (para continuar vMerge)
    active_col_master = {}
    for ri, tr in enumerate(trs):
        col = 0
        for tc in tr.findall(qn('w:tc')):
            tcPr = tc.find(qn('w:tcPr'))
            gs = 1
            vcont = False; vstart = False
            if tcPr is not None:
                gsEl = tcPr.find(qn('w:gridSpan'))
                if gsEl is not None:
                    try: gs = int(gsEl.get(qn('w:val'), '1'))
                    except Exception: pass
                vm = tcPr.find(qn('w:vMerge'))
                if vm is not None:
                    val = vm.get(qn('w:val'))
                    if val == 'restart': vstart = True
                    else: vcont = True  # continue (sem val ou val=continue)
            if vcont and col in active_col_master:
                # Estender rowspan do master aberto nesta coluna
                mi = active_col_master[col]
                masters[mi]["rowspan"] = ri - masters[mi]["r"] + 1
            else:
                masters.append({
                    "r": ri, "c": col,
                    "rowspan": 1, "colspan": gs,
                    "text": _cell_text(tc),
                })
                active_col_master[col] = len(masters) - 1
            col += gs
    return {"n_rows": n_rows, "n_cols": n_cols, "cells": masters}

def _extrair_tabelas_em_ordem(filepath):
    """Percorre o body do docx em ordem e extrai tabelas nao-gabarito,
    anotando qual foi a ultima questao vista antes de cada tabela.

    Retorna lista de dicts: [{"qnum_antes", "n_rows", "n_cols", "cells", "titulo"}]
    """
    from docx import Document
    doc = Document(filepath)
    body = doc.element.body
    RE_NUM = re.compile(r'^(\d{1,2})\.\s+')
    ultimo_qnum = 0
    tbl_idx = 0  # indice sequencial para casar com doc.tables
    saida = []
    for child in body.iterchildren():
        tag = child.tag
        if tag == qn('w:p'):
            # pega texto do paragrafo via w:t
            texts = child.findall('.//' + qn('w:t'))
            txt = ''.join((t.text or '') for t in texts).strip()
            m = RE_NUM.match(txt)
            if m:
                try: ultimo_qnum = int(m.group(1))
                except Exception: pass
        elif tag == qn('w:tbl'):
            if tbl_idx < len(doc.tables):
                tbl = doc.tables[tbl_idx]
                tbl_idx += 1
                if _tabela_eh_gabarito(tbl):
                    continue  # gabarito eh processado em outro lugar
                info = _extrair_tabela_como_dict(tbl)
                if info:
                    # Usar primeira celula como titulo se for clara
                    # (ex: "ALTERAÇÃO ECONÔMICA DO CONTRATO" mergeada em todo o header)
                    titulo = None
                    if info["cells"] and info["cells"][0]["r"] == 0 \
                       and info["cells"][0]["c"] == 0 \
                       and info["cells"][0]["colspan"] >= info["n_cols"]:
                        titulo = info["cells"][0]["text"]
                        # Remove a linha de titulo da tabela
                        info["cells"] = info["cells"][1:]
                        info["n_rows"] -= 1
                        for c in info["cells"]:
                            c["r"] -= 1
                    saida.append({
                        "qnum_antes": ultimo_qnum,
                        "n_rows": info["n_rows"],
                        "n_cols": info["n_cols"],
                        "cells": info["cells"],
                        "titulo": titulo,
                    })
    return saida

def _inserir_slides_tabela(slides, tabelas_info):
    """Insere slides de tabela apos a questao correspondente."""
    if not tabelas_info:
        return slides
    por_q = {}
    for t in tabelas_info:
        por_q.setdefault(t["qnum_antes"], []).append(t)
    qnums = {s.get("numero") for s in slides if s.get("tipo") == "questao"}
    novos = []
    # Tabelas antes da primeira questao
    for k in list(por_q.keys()):
        if k == 0 or k not in qnums:
            for t in por_q.pop(k):
                novos.append({"tipo":"tabela", "n_rows":t["n_rows"], "n_cols":t["n_cols"],
                              "cells":t["cells"], "titulo":t.get("titulo")})
    for s in slides:
        novos.append(s)
        if s.get("tipo") == "questao":
            q = s.get("numero")
            for t in por_q.pop(q, []):
                novos.append({"tipo":"tabela", "n_rows":t["n_rows"], "n_cols":t["n_cols"],
                              "cells":t["cells"], "titulo":t.get("titulo")})
    # Restos
    for k, lst in por_q.items():
        for t in lst:
            novos.append({"tipo":"tabela", "n_rows":t["n_rows"], "n_cols":t["n_cols"],
                          "cells":t["cells"], "titulo":t.get("titulo")})
    return novos

def _extrair_texto_e_imgs_docx(filepath):
    """Extrai texto bruto + imagens do docx (body + tabelas).

    Para cada imagem, guarda o numero da ultima questao detectada antes dela
    (padrao 'N. ...' no inicio da linha).
    Retorna (texto, imgs) onde imgs = [{"qnum_antes": int, "b64": str, "ext": str}]
    """
    from docx import Document
    doc = Document(filepath)
    doc_part = doc.part
    linhas = []
    imgs = []
    RE_NUM_HEAD = re.compile(r'^(\d{1,2})[.\-)]\s+')
    ultimo_qnum = 0

    def _processar_para(para):
        nonlocal ultimo_qnum
        txt = para.text.strip()
        m = RE_NUM_HEAD.match(txt) if txt else None
        if m:
            try: ultimo_qnum = int(m.group(1))
            except Exception: pass
        if txt:
            linhas.append(txt)
        else:
            linhas.append("")
        for b64, ext in _get_imgs_from_para(para, doc_part):
            imgs.append({"qnum_antes": ultimo_qnum, "b64": b64, "ext": ext})

    # Paragrafos do corpo
    for para in doc.paragraphs:
        _processar_para(para)
    # Paragrafos dentro de celulas de tabela (gabaritos, enunciados em tabela, etc.)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _processar_para(para)

    return "\n".join(linhas), imgs

def _reinjetar_imagens_nos_slides(slides, imgs):
    """Insere imagens apos a questao de numero correspondente."""
    if not imgs:
        return slides
    por_q = {}
    for im in imgs:
        por_q.setdefault(im["qnum_antes"], []).append(im)
    qnums_existentes = {s.get("numero") for s in slides if s.get("tipo") == "questao"}
    novos = []
    orfas_keys = [k for k in por_q.keys() if k == 0 or k not in qnums_existentes]
    for k in orfas_keys:
        for im in por_q.pop(k, []):
            novos.append({"tipo":"imagem","img_b64":im["b64"],"img_ext":im["ext"]})
    for s in slides:
        novos.append(s)
        if s.get("tipo") == "questao":
            q = s.get("numero")
            for im in por_q.pop(q, []):
                novos.append({"tipo":"imagem","img_b64":im["b64"],"img_ext":im["ext"]})
    for k, lst in por_q.items():
        for im in lst:
            novos.append({"tipo":"imagem","img_b64":im["b64"],"img_ext":im["ext"]})
    return novos

def _parse_via_claude(texto_bruto, disciplina="", assunto=""):
    """Usa Claude API para extrair questões.

    Divide em blocos SEMPRE que o texto for consideravelmente grande,
    para manter cada chamada rapida (< 90s) e nao estourar timeout do
    servidor em arquivos grandes. Cada chamada tem retry automatico.
    """
    import time

    def _montar_prompt(bloco, parte_n=None, total_partes=None):
        sufixo = ""
        if total_partes and total_partes > 1:
            sufixo = f"\n\n(Este e o bloco {parte_n} de {total_partes}. Extraia apenas as questoes deste bloco.)"
        return f"""Voce e um extrator de questoes de concurso. Analise o texto abaixo e extraia TODAS as questoes.

Retorne APENAS um JSON valido, sem texto antes ou depois, sem marcacoes markdown, com esta estrutura exata:
{{
  "questoes": [
    {{
      "numero": 1,
      "enunciado": "texto completo do enunciado",
      "alternativas": ["A) texto", "B) texto", "C) texto", "D) texto", "E) texto"],
      "certo_errado": false
    }}
  ],
  "gabarito": {{
    "questoes": [1, 2, 3],
    "respostas": ["A", "B", "C"]
  }}
}}

Regras:
- Capture o enunciado COMPLETO incluindo textos de apoio/excertos antes das alternativas
- Alternativas sempre no formato "A) texto", "B) texto" etc
- Se nao houver gabarito no texto, retorne gabarito com listas vazias
- Se for questao Certo/Errado, coloque certo_errado=true e alternativas=[]
- Numere sequencialmente se nao houver numeracao explicita
- Processe TODAS as questoes do texto, incluindo a ultima. Nao pare no meio.{sufixo}

TEXTO:
{bloco}
"""

    def _chamar_com_retry(prompt, max_tokens=8000, timeout=90, tentativas=3):
        """Chama Claude API com retry exponencial em caso de erro de rede/timeout."""
        ultimo_erro = None
        for t in range(tentativas):
            try:
                return _chamar_claude_api(prompt, max_tokens=max_tokens, timeout=timeout)
            except Exception as e:
                ultimo_erro = e
                traceback.print_exc()
                if t < tentativas - 1:
                    time.sleep(2 ** t)  # 1s, 2s, 4s
        raise ultimo_erro if ultimo_erro else RuntimeError("Falha ao chamar Claude API")

    # Divide em blocos de ate 60k chars para manter cada chamada < 90s.
    # (Antes era 150k e dava timeout em docs grandes.)
    LIMITE_BLOCO = 60000

    # Texto pequeno: 1 chamada so.
    if len(texto_bruto) <= LIMITE_BLOCO:
        return _chamar_com_retry(_montar_prompt(texto_bruto), max_tokens=8000, timeout=90)

    # Texto grande: divide por marcadores de questao para nao cortar no meio.
    RE_Q = re.compile(r'(?=^\d{1,2}\.\s)', re.MULTILINE)
    pedacos = RE_Q.split(texto_bruto)
    blocos, atual = [], ""
    for p in pedacos:
        if len(atual) + len(p) > LIMITE_BLOCO and atual:
            blocos.append(atual); atual = p
        else:
            atual += p
    if atual: blocos.append(atual)

    all_q, all_gq, all_gr = [], [], []
    for i, b in enumerate(blocos):
        r = _chamar_com_retry(_montar_prompt(b, i+1, len(blocos)),
                              max_tokens=8000, timeout=90)
        all_q.extend(r.get("questoes", []))
        gab = r.get("gabarito", {}) or {}
        all_gq.extend(gab.get("questoes", []))
        all_gr.extend(gab.get("respostas", []))
    return {"questoes": all_q, "gabarito": {"questoes": all_gq, "respostas": all_gr}}

def _parse_docx(filepath):
    from docx import Document
    doc = Document(filepath)
    doc_part = doc.part
    slides, gqs, grs = [], [], []
    RE_NUM      = re.compile(r'^(\d{1,2})\.\s+(.+)', re.DOTALL)
    RE_NUM_ONLY = re.compile(r'^(\d{1,2})\.\s*')
    RE_LETRA    = re.compile(r'^[A-Ea-e]' + '$')
    RE_LETRA_SOLTA = re.compile(r'^([A-Ea-e])\s+(.+)', re.DOTALL)
    RE_ALT      = re.compile(r'^\([A-Ea-e]\)\s+.+|^[A-Ea-e][).]\s+.+')
    RE_CE       = re.compile(r'^(certo|errado)[\s.(]', re.IGNORECASE)
    RE_GAB_CELL = re.compile(r'^(\d{1,2})\.\s*([A-Ea-eCcEe])')
    RE_GAB_TX   = re.compile(r'(\d{1,2})\s*([A-Ea-e])\b')
    RE_ALT_SEP  = re.compile(r'^Alternativas$', re.IGNORECASE)
    def ib(p): return any(r.bold for r in p.runs if r.text.strip())
    def hn(p): return p._element.find('.//' + WS_NS + 'numPr') is not None
    def has_img(p): return _has_img_xml(p._element.xml)
    # Gabarito em tabela
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                txt = cell.text.strip()
                m = RE_GAB_CELL.match(txt)
                if m: gqs.append(int(m.group(1))); grs.append(m.group(2).upper())
                else:
                    for mm in RE_GAB_TX.finditer(txt):
                        gqs.append(int(mm.group(1))); grs.append(mm.group(2).upper())
    # Agrupar em blocos por linha vazia
    blocos = []
    bloco = []
    for para in doc.paragraphs:
        txt = para.text.strip(); is_img = has_img(para)
        if not txt and not is_img:
            if bloco: blocos.append(bloco); bloco = []
        else: bloco.append(para)
    if bloco: blocos.append(bloco)

    # --- DETECÇÃO DE PADRÃO ---

    # Padrão 4: "Alternativas" como separador
    tem_alt_sep = sum(1 for p in doc.paragraphs if RE_ALT_SEP.match(p.text.strip()))
    if tem_alt_sep >= 2:
        qnum = 0
        state_num = None; state_enc = []; state_in_alts = False; state_alts = []
        def flush(slides, state_num, state_enc, state_alts, qnum):
            enc = ' '.join(e for e in state_enc if e)
            if state_num: qnum = state_num
            else: qnum += 1
            alts = [a['letra'] + ') ' + (a['texto'] or '') for a in state_alts if a.get('texto')]
            slides.append({'tipo':'questao','numero':qnum,'enunciado':enc,'certo_errado':False,'alternativas':alts})
            return qnum
        for para in doc.paragraphs:
            txt = para.text.strip()
            if has_img(para):
                if state_in_alts or state_enc:
                    qnum = flush(slides, state_num, state_enc, state_alts, qnum)
                    state_num = None; state_enc = []; state_in_alts = False; state_alts = []
                for ir in _get_imgs_from_para(para, doc_part):
                    slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                continue
            if not txt: continue
            m = RE_NUM_ONLY.match(txt)
            if m:
                if state_in_alts or state_enc:
                    qnum = flush(slides, state_num, state_enc, state_alts, qnum)
                state_num = int(m.group(1))
                state_enc = [txt]; state_in_alts = False; state_alts = []
                continue
            if RE_ALT_SEP.match(txt):
                state_in_alts = True; continue
            if state_in_alts:
                if RE_LETRA.match(txt): state_alts.append({'letra': txt, 'texto': None})
                elif state_alts and state_alts[-1]['texto'] is None: state_alts[-1]['texto'] = txt
                continue
            if state_num is not None: state_enc.append(txt)
        if state_in_alts or state_enc:
            flush(slides, state_num, state_enc, state_alts, qnum)
        gab = {"questoes": gqs, "respostas": grs} if gqs else None
        return slides, gab

    # Padrão 5: alternativas "A texto", "B texto" em blocos separados por linha vazia
    def is_alt_block(bloco):
        txts = [p.text.strip() for p in bloco]
        if len(txts) < 4: return False
        letras = []
        for t in txts:
            m = RE_LETRA_SOLTA.match(t)
            if m: letras.append(m.group(1).upper())
            else: return False
        return letras == list('ABCDE')[:len(letras)]

    tem_alt_blocos = sum(1 for b in blocos if is_alt_block(b))
    if tem_alt_blocos >= 2:
        qnum = 0; i = 0
        while i < len(blocos):
            bloco = blocos[i]
            txts = [p.text.strip() for p in bloco]
            imgs_bloco = []
            for p in bloco:
                if has_img(p):
                    imgs_bloco.extend(_get_imgs_from_para(p, doc_part))
            if any('gabarito' in t.lower() for t in txts):
                for t in txts:
                    for m in RE_GAB_TX.finditer(t):
                        gqs.append(int(m.group(1))); grs.append(m.group(2).upper())
                i += 1; continue
            if is_alt_block(bloco):
                for ir in imgs_bloco:
                    slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                i += 1; continue
            enc = ' '.join(t for t in txts if t)
            if i + 1 < len(blocos) and is_alt_block(blocos[i+1]):
                qnum += 1
                alts = []
                for t in [p.text.strip() for p in blocos[i+1]]:
                    m = RE_LETRA_SOLTA.match(t)
                    if m: alts.append(m.group(1).upper() + ') ' + m.group(2).strip())
                slides.append({'tipo':'questao','numero':qnum,'enunciado':enc,'certo_errado':False,'alternativas':alts})
                for ir in imgs_bloco:
                    slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                i += 2
            else:
                if enc:
                    slides.append({'tipo':'contexto','texto':enc})
                for ir in imgs_bloco:
                    slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                i += 1
        gab = {"questoes": gqs, "respostas": grs} if gqs else None
        return slides, gab

    # Padrão 3: blocos por linha vazia com alternativas (A), (B)...
    tem_alts_par = sum(1 for b in blocos
                       if any(p.text.strip().startswith('(') and RE_ALT.match(p.text.strip()) for p in b))
    if tem_alts_par >= 2:
        qnum = 0
        for bloco in blocos:
            txts = [p.text.strip() for p in bloco]
            imgs_bloco = []
            for p in bloco:
                if has_img(p):
                    imgs_bloco.extend(_get_imgs_from_para(p, doc_part))
            if any(t.upper() == 'GABARITO' for t in txts):
                for t in txts:
                    for m in RE_GAB_TX.finditer(t):
                        gqs.append(int(m.group(1))); grs.append(m.group(2).upper())
                continue
            alts_idx = [i for i, t in enumerate(txts) if RE_ALT.match(t)]
            if not alts_idx:
                txt_c = ' '.join(t for t in txts if t)
                if txt_c: slides.append({'tipo':'contexto','texto':txt_c})
                for ir in imgs_bloco:
                    slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                continue
            first_alt = alts_idx[0]
            enc_txt = ' '.join(txts[:first_alt]) if txts[:first_alt] else (txts[0] if txts else '')
            alts = [txts[j] for j in alts_idx]
            m_num = RE_NUM.match(enc_txt) if enc_txt else None
            if m_num: qnum = int(m_num.group(1))
            else: qnum += 1
            ce = len(alts) > 0 and all(RE_CE.match(a) for a in alts)
            slides.append({'tipo':'questao','numero':qnum,'enunciado':enc_txt,'certo_errado':ce,'alternativas':alts if not ce else []})
            for ir in imgs_bloco:
                slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
        gab = {"questoes": gqs, "respostas": grs} if gqs else None
        return slides, gab

    # Padrão 6 (CEBRASPE): questões numeradas "01. texto" sem alternativas, certo/errado
    # Gabarito na tabela como "01.V" ou "01.F"
    RE_GAB_CE = re.compile(r'^(\d{1,2})\.\s*([VFvf])\b')
    # Coletar gabarito V/F da tabela
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                for mm in RE_GAB_CE.finditer(cell.text.strip()):
                    q_n = int(mm.group(1))
                    letra = mm.group(2).upper()
                    if q_n not in gqs:  # evitar duplicata com RE_GAB_CELL
                        gqs.append(q_n)
                        grs.append('C' if letra == 'V' else 'E')
    tem_ce_tabela = any(RE_GAB_CE.search(c.text) for tbl in doc.tables for row in tbl.rows for c in row.cells)
    tem_num_seq = sum(1 for p in doc.paragraphs if RE_NUM_ONLY.match(p.text.strip()))
    nao_tem_alts = not (tem_alt_sep >= 2 or tem_alts_par >= 2 or tem_alt_blocos >= 2)
    usa_cebraspe = tem_ce_tabela and tem_num_seq >= 2 and nao_tem_alts
    if usa_cebraspe:
        state_enc = []; state_qnum = None
        def flush_ce(state_qnum, state_enc, slides):
            if state_qnum and state_enc:
                enc = ' '.join(e for e in state_enc if e)
                slides.append({'tipo':'questao','numero':state_qnum,'enunciado':enc,'certo_errado':True,'alternativas':[]})
        for para in doc.paragraphs:
            txt = para.text.strip()
            if has_img(para):
                flush_ce(state_qnum, state_enc, slides)
                state_qnum = None; state_enc = []
                for ir in _get_imgs_from_para(para, doc_part):
                    slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                continue
            if not txt or txt.upper() == 'GABARITO': continue
            m = RE_NUM_ONLY.match(txt)
            if m:
                flush_ce(state_qnum, state_enc, slides)
                state_qnum = int(m.group(1)); state_enc = [txt]
            elif state_qnum is not None:
                state_enc.append(txt)
        flush_ce(state_qnum, state_enc, slides)
        gab = {"questoes": gqs, "respostas": grs} if gqs else None
        return slides, gab

    # Padrões 1 e 2: bold / List Paragraph
    paras = doc.paragraphs; i, qnum = 0, 0
    while i < len(paras):
        para = paras[i]; txt = para.text.strip()
        if has_img(para):
            for ir in _get_imgs_from_para(para, doc_part):
                slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
            i += 1; continue
        if not txt: i += 1; continue
        bold = ib(para); m = RE_NUM.match(txt)
        if m and bold:
            qnum = int(m.group(1)); enc = txt; i += 1; alts = []; extras = []
            while i < len(paras):
                p2 = paras[i]; t2 = p2.text.strip()
                if has_img(p2):
                    for ir in _get_imgs_from_para(p2, doc_part):
                        slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                    i += 1; continue
                if not t2: i += 1; continue
                b2 = ib(p2); m2 = RE_NUM.match(t2)
                if m2 and b2: break
                if p2.style.name == 'List Paragraph' and hn(p2): break
                if RE_CE.match(t2): alts.append(t2); i += 1
                elif RE_ALT.match(t2): alts.append(t2); i += 1
                elif b2 and RE_LETRA.match(t2):
                    letra = t2; i += 1
                    if i < len(paras) and paras[i].text.strip():
                        alts.append(letra + ') ' + paras[i].text.strip()); i += 1
                else: extras.append(t2); i += 1
            enc2 = enc + ('\n' + '\n'.join(extras) if extras else '')
            ce = len(alts) > 0 and all(RE_CE.match(a) for a in alts)
            slides.append({'tipo':'questao','numero':qnum,'enunciado':enc2,'certo_errado':ce,'alternativas':alts if not ce else []})
            continue
        if para.style.name == 'List Paragraph' and hn(para) and txt:
            qnum += 1; enc = txt; i += 1; alts = []; extras = []
            while i < len(paras):
                p2 = paras[i]; t2 = p2.text.strip()
                if has_img(p2):
                    for ir in _get_imgs_from_para(p2, doc_part):
                        slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
                    i += 1; continue
                if not t2: i += 1; continue
                b2 = ib(p2)
                if p2.style.name == 'List Paragraph' and hn(p2): break
                if RE_NUM.match(t2) and b2: break
                if b2 and RE_LETRA.match(t2):
                    letra = t2; i += 1
                    if i < len(paras):
                        t3 = paras[i].text.strip(); b3 = ib(paras[i])
                        if t3 and not (b3 and RE_LETRA.match(t3)):
                            alts.append(letra + ') ' + t3); i += 1
                    continue
                if RE_CE.match(t2): alts.append(t2); i += 1; continue
                if RE_ALT.match(t2): alts.append(t2); i += 1; continue
                if not alts: extras.append(t2)
                i += 1
            enc2 = enc + ('\n' + '\n'.join(extras) if extras else '')
            ce = len(alts) > 0 and all(RE_CE.match(a) for a in alts)
            slides.append({'tipo':'questao','numero':qnum,'enunciado':enc2,'certo_errado':ce,'alternativas':alts if not ce else []})
            continue
        i += 1
    # Varredura extra: imagens dentro de celulas de tabela (nao aparecem em doc.paragraphs)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                for p_cell in cell.paragraphs:
                    if has_img(p_cell):
                        for ir in _get_imgs_from_para(p_cell, doc_part):
                            slides.append({'tipo':'imagem','img_b64':ir[0],'img_ext':ir[1]})
    gab = {"questoes": gqs, "respostas": grs} if gqs else None
    return slides, gab

def _parse_texto(texto):
    linhas = [l.rstrip() for l in texto.splitlines()]
    slides, gqs, grs = [], [], []
    RE_Q = re.compile(r'^(\d{1,2})[.\-]\s+(.+)')
    RE_A = re.compile(r'^[A-Ea-e][).]')
    RE_G = re.compile(r'^GABARITO', re.IGNORECASE)
    RE_GI = re.compile(r'(\d{1,2})\s*[-]\s*([A-Ea-eCcEe])\b')
    i = 0
    while i < len(linhas):
        l = linhas[i].strip()
        if RE_G.match(l):
            bloco = l; i += 1
            while i < len(linhas): bloco += ' ' + linhas[i].strip(); i += 1
            for m in RE_GI.finditer(bloco): gqs.append(int(m.group(1))); grs.append(m.group(2).upper())
            continue
        m = RE_Q.match(l)
        if m:
            num = int(m.group(1)); ep = [m.group(2).strip()]; i += 1; alts = []; ce = False
            while i < len(linhas):
                ll = linhas[i].strip()
                if not ll:
                    i += 1
                    if i < len(linhas) and (RE_Q.match(linhas[i].strip()) or RE_G.match(linhas[i].strip())): break
                    continue
                if RE_A.match(ll): alts.append(ll); i += 1
                elif ll.lower().startswith('certo') or ll.lower().startswith('errado'): ce = True; i += 1
                elif RE_Q.match(ll) or RE_G.match(ll): break
                else: ep.append(ll); i += 1
            slides.append({"tipo":"questao","numero":num,"enunciado":" ".join(ep),"certo_errado":ce,"alternativas":alts})
            continue
        if l and not RE_G.match(l):
            cp = [l]; i += 1
            while i < len(linhas):
                ll = linhas[i].strip()
                if RE_Q.match(ll) or RE_G.match(ll): break
                if not ll:
                    i += 1
                    if i < len(linhas) and not linhas[i].strip(): break
                    continue
                cp.append(ll); i += 1
            slides.append({"tipo":"contexto","texto":" ".join(cp)})
            continue
        i += 1
    gab = {"questoes": gqs, "respostas": grs} if gqs else None
    return slides, gab

def _parse_pptx(filepath):
    """Extrai texto e imagens de um PPTX não formatado, preservando estrutura slide a slide."""
    prs_in = Presentation(filepath)
    slides_parsed = []
    all_text_lines = []  # para tentativa de parse de questões

    # --- Passo 1: extrair estrutura de cada slide ---
    extracted = []  # lista de dicts: {titulo, subtitulo, conteudos, imagens}
    for slide in prs_in.slides:
        info = {"titulo": "", "subtitulo": "", "conteudos": [], "imagens": []}
        for shape in slide.shapes:
            # Imagens
            if shape.shape_type == 13:
                try:
                    blob = shape.image.blob
                    ct = shape.image.content_type or "image/png"
                    ext = ct.split("/")[-1].replace("jpeg","jpg")
                    if ext not in ("png","jpg","gif","bmp","webp"): ext = "png"
                    info["imagens"].append((base64.b64encode(blob).decode(), ext))
                except Exception:
                    pass
            # Texto
            if shape.has_text_frame:
                name = shape.name.lower()
                texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                if 'título' in name or 'title' in name:
                    info["titulo"] = " ".join(texts)
                elif 'subtítulo' in name or 'subtitle' in name:
                    info["subtitulo"] = " ".join(texts)
                else:
                    info["conteudos"].extend(texts)
            # Tabelas
            if shape.has_table:
                for row in shape.table.rows:
                    row_t = [c.text.strip() for c in row.cells if c.text.strip()]
                    if row_t:
                        info["conteudos"].append(" | ".join(row_t))
        extracted.append(info)
        # Acumular texto para tentativa de parse de questões
        if info["titulo"]: all_text_lines.append(info["titulo"])
        for c in info["conteudos"]: all_text_lines.append(c)
        all_text_lines.append("")

    # --- Passo 2: tentar detectar se é material de questões ---
    texto_bruto = "\n".join(all_text_lines)
    questoes_parsed, gab = _parse_texto(texto_bruto)
    tem_questoes = any(s.get("tipo") == "questao" for s in questoes_parsed)

    if tem_questoes:
        # É material de questões — usar parse de questões + imagens soltas
        for info in extracted:
            for img_b64, img_ext in info["imagens"]:
                questoes_parsed.append({"tipo": "imagem", "img_b64": img_b64, "img_ext": img_ext})
        return questoes_parsed, gab

    # --- Passo 3: não é questões — preservar estrutura slide a slide ---
    for info in extracted:
        titulo = info["titulo"]
        subtitulo = info["subtitulo"]
        conteudos = info["conteudos"]
        imagens = info["imagens"]
        is_divider = not conteudos and not imagens and (titulo or subtitulo)

        if is_divider:
            # Slide divisor de seção
            label = titulo
            if subtitulo:
                label = (label + " — " + subtitulo) if label else subtitulo
            slides_parsed.append({"tipo": "secao", "titulo": label})
        else:
            # Slide de conteúdo com título
            slides_parsed.append({
                "tipo": "conteudo_slide",
                "titulo": titulo,
                "paragrafos": conteudos
            })
            # Imagens do slide
            for img_b64, img_ext in imagens:
                slides_parsed.append({"tipo": "imagem", "img_b64": img_b64, "img_ext": img_ext})

    return slides_parsed, None

def _extrair_conteudo_pptx(filepath):
    """Extrai texto estruturado (slide a slide) e imagens de um PPTX."""
    prs_in = Presentation(filepath)
    slides_info = []
    img_list = []  # (slide_index, b64, ext)

    for si, slide in enumerate(prs_in.slides):
        titulo = ""
        subtitulo = ""
        conteudos = []
        for shape in slide.shapes:
            if shape.shape_type == 13:
                try:
                    blob = shape.image.blob
                    ct = shape.image.content_type or "image/png"
                    ext = ct.split("/")[-1].replace("jpeg","jpg")
                    if ext not in ("png","jpg","gif","bmp","webp"): ext = "png"
                    img_list.append((si, base64.b64encode(blob).decode(), ext))
                except Exception:
                    pass
            if shape.has_text_frame:
                name = shape.name.lower()
                texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                if 'título' in name or 'title' in name:
                    titulo = " ".join(texts)
                elif 'subtítulo' in name or 'subtitle' in name:
                    subtitulo = " ".join(texts)
                else:
                    conteudos.extend(texts)
            if shape.has_table:
                for row in shape.table.rows:
                    row_t = [c.text.strip() for c in row.cells if c.text.strip()]
                    if row_t: conteudos.append(" | ".join(row_t))
        slides_info.append({"index": si, "titulo": titulo, "subtitulo": subtitulo, "conteudos": conteudos})

    # Montar texto estruturado para enviar à IA
    texto_estruturado = ""
    for s in slides_info:
        texto_estruturado += f"\n--- SLIDE {s['index']+1} ---\n"
        if s["titulo"]: texto_estruturado += f"TÍTULO: {s['titulo']}\n"
        if s["subtitulo"]: texto_estruturado += f"SUBTÍTULO: {s['subtitulo']}\n"
        for c in s["conteudos"]:
            texto_estruturado += f"{c}\n"

    return texto_estruturado, img_list

_CLAUDE_PROMPT_SLIDES = """Você é um especialista em formatação de materiais educativos para concursos públicos.

Analise o conteúdo extraído de uma apresentação PowerPoint abaixo. Este material pode conter:
- Slides de TEORIA (artigos de lei, jurisprudência, doutrina, explicações)
- Slides de QUESTÕES (enunciados com alternativas A-E ou Certo/Errado)
- Slides DIVISORES de seção (apenas título e subtítulo, sem conteúdo)
- Qualquer combinação dos anteriores

Sua tarefa é classificar e estruturar CADA slide corretamente.

Retorne APENAS um JSON válido, sem texto antes ou depois, sem marcações markdown, com esta estrutura:
{
  "slides": [
    {
      "tipo": "secao",
      "titulo": "Texto do título da seção"
    },
    {
      "tipo": "conteudo_slide",
      "titulo": "Título do slide",
      "paragrafos": ["parágrafo 1", "parágrafo 2"]
    },
    {
      "tipo": "questao",
      "numero": 1,
      "enunciado": "texto completo do enunciado incluindo textos de apoio",
      "alternativas": ["A) texto", "B) texto", "C) texto", "D) texto", "E) texto"],
      "certo_errado": false
    }
  ],
  "gabarito": {
    "questoes": [1, 2],
    "respostas": ["A", "B"]
  }
}

REGRAS IMPORTANTES:
1. Slides que têm APENAS título (e talvez subtítulo) sem conteúdo são "secao". Combine título e subtítulo: "Título — Subtítulo"
2. Slides com título + parágrafos de teoria/lei/jurisprudência são "conteudo_slide". Mantenha CADA parágrafo separado no array.
3. Slides com questões numeradas + alternativas são "questao". Alternativas no formato "A) texto".
4. Se a questão for Certo/Errado (sem alternativas A-E), coloque certo_errado=true e alternativas=[]
5. PRESERVE o conteúdo completo de cada parágrafo, sem resumir nem cortar.
6. Se não houver gabarito, retorne gabarito com listas vazias.
7. Mantenha a ORDEM original dos slides.
"""

def _chamar_claude_api(prompt_text, max_tokens=8000, timeout=90):
    """Chamada generica a Claude API. Retorna o JSON da resposta parseado.

    Levanta excecao especifica para erros HTTP (incluindo 429/529) para
    que o caller possa decidir se tenta novamente.
    """
    import urllib.request, urllib.error, json as jsonlib
    payload = jsonlib.dumps({
        "model": "claude-sonnet-4-20250514",
        "max_tokens": max_tokens,
        "messages": [{"role": "user", "content": prompt_text}]
    }).encode()
    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=payload,
        headers={
            "Content-Type": "application/json",
            "anthropic-version": "2023-06-01",
            "x-api-key": os.environ.get("ANTHROPIC_API_KEY", "")
        },
        method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            data = jsonlib.loads(resp.read())
    except urllib.error.HTTPError as he:
        # Repassar erro claro — retry sera feito pelo caller se aplicavel
        try:
            body = he.read().decode("utf-8", errors="ignore")
        except Exception:
            body = ""
        raise RuntimeError(f"Claude API HTTP {he.code}: {body[:300]}") from he
    except urllib.error.URLError as ue:
        raise RuntimeError(f"Claude API network error: {ue}") from ue
    except Exception as e:
        raise RuntimeError(f"Claude API error: {e}") from e
    if "content" not in data or not data["content"]:
        raise RuntimeError(f"Claude API resposta inesperada: {str(data)[:300]}")
    txt = data["content"][0]["text"].strip()
    if txt.startswith("```"): txt = "\n".join(txt.split("\n")[1:])
    if txt.endswith("```"): txt = "\n".join(txt.split("\n")[:-1])
    return jsonlib.loads(txt.strip())

def _dividir_texto_em_blocos(texto_estruturado, max_chars=10000):
    """Divide texto estruturado em blocos respeitando limites de slide."""
    partes = texto_estruturado.split("\n--- SLIDE ")
    blocos = []
    bloco_atual = ""
    for i, parte in enumerate(partes):
        pedaco = ("--- SLIDE " + parte) if i > 0 else parte
        if len(bloco_atual) + len(pedaco) > max_chars and bloco_atual:
            blocos.append(bloco_atual)
            bloco_atual = pedaco
        else:
            bloco_atual += ("\n" if bloco_atual else "") + pedaco
    if bloco_atual:
        blocos.append(bloco_atual)
    return blocos

def _parse_pptx_via_claude(filepath, disciplina="", assunto=""):
    """Usa Claude API para analisar PPTX e estruturar conteúdo inteligentemente."""
    import json as jsonlib
    import time

    texto_estruturado, img_list = _extrair_conteudo_pptx(filepath)

    # Blocos menores (8k chars) para cada chamada rodar em <90s
    blocos = _dividir_texto_em_blocos(texto_estruturado, max_chars=8000)

    def _chamar_com_retry(prompt, tentativas=3):
        ultimo_erro = None
        for t in range(tentativas):
            try:
                return _chamar_claude_api(prompt, max_tokens=8000, timeout=90)
            except Exception as e:
                ultimo_erro = e
                traceback.print_exc()
                if t < tentativas - 1:
                    time.sleep(2 ** t)
        raise ultimo_erro if ultimo_erro else RuntimeError("Falha Claude API")

    all_slides = []
    all_gab_q = []
    all_gab_r = []

    for bi, bloco in enumerate(blocos):
        sufixo = ""
        if len(blocos) > 1:
            sufixo = f"\n\n(Este é o bloco {bi+1} de {len(blocos)}. Processe apenas os slides deste bloco.)"

        prompt = _CLAUDE_PROMPT_SLIDES + f"\nCONTEÚDO DA APRESENTAÇÃO:{sufixo}\n{bloco}"
        resultado = _chamar_com_retry(prompt)

        all_slides.extend(resultado.get("slides", []))
        gab = resultado.get("gabarito", {})
        if gab and gab.get("questoes"):
            all_gab_q.extend(gab["questoes"])
            all_gab_r.extend(gab["respostas"])

    # Adicionar imagens extraídas
    for si, img_b64, img_ext in img_list:
        all_slides.append({"tipo": "imagem", "img_b64": img_b64, "img_ext": img_ext})

    gab_final = {"questoes": all_gab_q, "respostas": all_gab_r} if all_gab_q else None
    return all_slides, gab_final


app = Flask(__name__)
CORS(app, origins="*")

@app.route("/", methods=["GET"])
def health():
    return jsonify({"status": "ok"})

@app.route("/gerar", methods=["POST"])
def gerar():
    try:
        if request.files or request.form:
            arq = request.files.get("arquivo")
            disc = request.form.get("disciplina","DISCIPLINA")
            ass  = request.form.get("assunto","ASSUNTO")
            prof = request.form.get("professor","")
            tipo = request.form.get("tipo","QUESTOES")
            formato = request.form.get("formato","word_to_slides")
            usar_ia = request.form.get("usar_ia","0") == "1"
            if not arq: return jsonify({"erro":"Arquivo nao enviado"}), 400
            fname = arq.filename.lower()

            # Auto-detectar formato pela extensão se necessário
            is_pptx = fname.endswith(".pptx")
            is_docx = fname.endswith(".docx")

            # Flag: saida em Word em vez de PowerPoint
            saida_word = formato in ("slides_to_word", "word_to_word")

            # --- SLIDES → WORD ---
            if formato == "slides_to_word":
                if not is_pptx:
                    return jsonify({"erro":"Para o formato Slides→Word, envie um arquivo .pptx"}), 400
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                    arq.save(tmp.name); path = tmp.name
                try:
                    if usar_ia:
                        sl, gab = _parse_pptx_via_claude(path, disc, ass)
                    else:
                        sl, gab = _parse_pptx(path)
                finally: os.unlink(path)
                payload = {"disciplina":disc,"assunto":ass,"tipo":tipo,"professor":prof,"slides":sl,"gabarito":gab}
                buf = _gerar_docx(payload)
                d = disc.replace(" ","_") or "apresentacao"
                a = ass.replace(" ","_")
                fn = "Carranza_" + d + "_" + a + ".docx" if a else "Carranza_" + d + ".docx"
                return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                 as_attachment=True, download_name=fn)

            # --- WORD → WORD ---
            # Compartilha TODO o pipeline de extracao do word_to_slides
            # (parse docx + tabelas + imagens + IA opcional) e no final
            # gera .docx formatado no lugar de .pptx.
            elif formato == "word_to_word":
                if not is_docx:
                    return jsonify({"erro":"Para o formato Word→Word, envie um arquivo .docx"}), 400
                with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                    arq.save(tmp.name); path = tmp.name
                try:
                    if usar_ia:
                        texto_bruto, imgs_docx = _extrair_texto_e_imgs_docx(path)
                        resultado = _parse_via_claude(texto_bruto, disc, ass)
                        if "slides" in resultado:
                            sl = resultado["slides"]
                        else:
                            sl = []
                            for q in resultado.get("questoes", []):
                                sl.append({
                                    "tipo": "questao",
                                    "numero": q.get("numero", 0),
                                    "enunciado": q.get("enunciado", ""),
                                    "certo_errado": q.get("certo_errado", False),
                                    "alternativas": q.get("alternativas", [])
                                })
                        gab_raw = resultado.get("gabarito", {})
                        gab = gab_raw if gab_raw and gab_raw.get("questoes") else None
                        sl = _reinjetar_imagens_nos_slides(sl, imgs_docx)
                    else:
                        sl, gab = _parse_docx(path)
                    try:
                        tabelas = _extrair_tabelas_em_ordem(path)
                        sl = _inserir_slides_tabela(sl, tabelas)
                    except Exception:
                        traceback.print_exc()
                finally: os.unlink(path)

            # --- SLIDES → SLIDES ---
            elif formato == "slides_to_slides" or is_pptx:
                if not is_pptx:
                    return jsonify({"erro":"Para o formato Slides→Slides, envie um arquivo .pptx"}), 400
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                    arq.save(tmp.name); path = tmp.name
                try:
                    if usar_ia:
                        sl, gab = _parse_pptx_via_claude(path, disc, ass)
                    else:
                        sl, gab = _parse_pptx(path)
                finally: os.unlink(path)

            elif is_docx:
                # --- WORD → SLIDES (fluxo original) ---
                with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                    arq.save(tmp.name); path = tmp.name
                try:
                    if usar_ia:
                        texto_bruto, imgs_docx = _extrair_texto_e_imgs_docx(path)
                        resultado = _parse_via_claude(texto_bruto, disc, ass)
                        if "slides" in resultado:
                            sl = resultado["slides"]
                        else:
                            sl = []
                            for q in resultado.get("questoes", []):
                                sl.append({
                                    "tipo": "questao",
                                    "numero": q.get("numero", 0),
                                    "enunciado": q.get("enunciado", ""),
                                    "certo_errado": q.get("certo_errado", False),
                                    "alternativas": q.get("alternativas", [])
                                })
                        gab_raw = resultado.get("gabarito", {})
                        gab = gab_raw if gab_raw and gab_raw.get("questoes") else None
                        # Re-injetar imagens que o Claude nao viu (ele so recebeu texto)
                        sl = _reinjetar_imagens_nos_slides(sl, imgs_docx)
                    else:
                        sl, gab = _parse_docx(path)
                    # Extrair tabelas de conteudo (nao-gabarito) e posicionar no fluxo
                    try:
                        tabelas = _extrair_tabelas_em_ordem(path)
                        sl = _inserir_slides_tabela(sl, tabelas)
                    except Exception:
                        traceback.print_exc()
                finally: os.unlink(path)
            else:
                # --- TXT → SLIDES (fluxo original) ---
                texto = arq.read().decode("utf-8", errors="ignore")
                sl, gab = _parse_texto(texto)

            payload = {"disciplina":disc,"assunto":ass,"tipo":tipo,"professor":prof,"slides":sl,"gabarito":gab}
        else:
            payload = request.get_json(force=True)
            if not payload: return jsonify({"erro":"Payload vazio"}), 400
            saida_word = False

        # --- Gerar saida: Word (.docx) ou PowerPoint (.pptx) ---
        d = payload.get("disciplina","apresentacao").replace(" ","_")
        a = payload.get("assunto","").replace(" ","_")
        if saida_word:
            buf = _gerar_docx(payload)
            fn = "Carranza_" + d + "_" + a + ".docx" if a else "Carranza_" + d + ".docx"
            return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                             as_attachment=True, download_name=fn)
        buf = _build_pptx(payload)
        fn = "Carranza_" + d + "_" + a + ".pptx" if a else "Carranza_" + d + ".pptx"
        return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                         as_attachment=True, download_name=fn)
    except Exception as e:
        traceback.print_exc()
        return jsonify({"erro": str(e)}), 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT",5000)), debug=False)
